"""
Document Processors

Handles document parsing and text extraction from various formats:
- PDF, DOCX, PPTX, TXT
- Metadata extraction
- Content cleaning and normalization
"""

import os
import io
from abc import ABC, abstractmethod
from typing import Optional, Dict, Any, List
from dataclasses import dataclass, field
from pathlib import Path
import hashlib

from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation


@dataclass
class Document:
    """Represents a processed document with text and metadata."""
    id: str
    content: str
    metadata: Dict[str, Any] = field(default_factory=dict)
    source_file: Optional[str] = None

    def to_dict(self) -> Dict[str, Any]:
        return {
            'id': self.id,
            'content': self.content,
            'metadata': self.metadata,
            'source_file': self.source_file,
        }


class DocumentProcessor(ABC):
    """Abstract base class for document processors."""

    @abstractmethod
    def can_process(self, file_path: str) -> bool:
        """Check if this processor can handle the given file."""
        pass

    @abstractmethod
    def process(self, file_path: str, metadata: Optional[Dict[str, Any]] = None) -> Document:
        """Process the file and return a Document object."""
        pass

    def _generate_id(self, file_path: str) -> str:
        """Generate a unique document ID based on file hash."""
        with open(file_path, 'rb') as f:
            file_hash = hashlib.md5(f.read()).hexdigest()
        return f"doc_{file_hash[:16]}"

    def _extract_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract basic metadata from file."""
        path = Path(file_path)
        stat = path.stat()

        return {
            'filename': path.name,
            'file_extension': path.suffix,
            'file_size': stat.st_size,
            'created_at': stat.st_ctime,
            'modified_at': stat.st_mtime,
        }


class PDFProcessor(DocumentProcessor):
    """Processor for PDF documents using pypdf."""

    def can_process(self, file_path: str) -> bool:
        return Path(file_path).suffix.lower() == '.pdf'

    def process(self, file_path: str, metadata: Optional[Dict[str, Any]] = None) -> Document:
        base_metadata = self._extract_metadata(file_path)
        if metadata:
            base_metadata.update(metadata)

        try:
            reader = PdfReader(file_path)

            # Extract text from all pages
            content_parts = []
            for page_num, page in enumerate(reader.pages):
                try:
                    text = page.extract_text()
                    if text.strip():
                        content_parts.append(text)
                except Exception as e:
                    content_parts.append(f"[Error extracting page {page_num}: {str(e)}]")

            content = '\n\n'.join(content_parts)

            # Add PDF-specific metadata
            base_metadata.update({
                'total_pages': len(reader.pages),
                'document_type': 'pdf',
                'is_encrypted': reader.is_encrypted,
            })

            # Extract PDF metadata if available
            if reader.metadata:
                pdf_metadata = {
                    'title': reader.metadata.get('/Title', ''),
                    'author': reader.metadata.get('/Author', ''),
                    'subject': reader.metadata.get('/Subject', ''),
                    'creator': reader.metadata.get('/Creator', ''),
                    'producer': reader.metadata.get('/Producer', ''),
                }
                base_metadata.update({
                    k: v for k, v in pdf_metadata.items() if v
                })

            return Document(
                id=self._generate_id(file_path),
                content=content,
                metadata=base_metadata,
                source_file=file_path,
            )

        except Exception as e:
            raise ValueError(f"Failed to process PDF: {str(e)}")


class DocxProcessor(DocumentProcessor):
    """Processor for DOCX documents using python-docx."""

    def can_process(self, file_path: str) -> bool:
        return Path(file_path).suffix.lower() in ('.docx', '.doc')

    def process(self, file_path: str, metadata: Optional[Dict[str, Any]] = None) -> Document:
        base_metadata = self._extract_metadata(file_path)
        if metadata:
            base_metadata.update(metadata)

        try:
            doc = DocxDocument(file_path)

            # Extract text from paragraphs
            content_parts = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content_parts.append(paragraph.text)

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        content_parts.append(row_text)

            content = '\n\n'.join(content_parts)

            # Add DOCX-specific metadata
            base_metadata.update({
                'total_paragraphs': len(doc.paragraphs),
                'total_tables': len(doc.tables),
                'document_type': 'docx',
            })

            # Extract core properties if available
            if doc.core_properties:
                core_props = {
                    'title': doc.core_properties.title or '',
                    'author': doc.core_properties.author or '',
                    'subject': doc.core_properties.subject or '',
                    'created': doc.core_properties.created,
                    'modified': doc.core_properties.modified,
                }
                base_metadata.update({
                    k: v for k, v in core_props.items() if v
                })

            return Document(
                id=self._generate_id(file_path),
                content=content,
                metadata=base_metadata,
                source_file=file_path,
            )

        except Exception as e:
            raise ValueError(f"Failed to process DOCX: {str(e)}")


class PPTXProcessor(DocumentProcessor):
    """Processor for PPTX presentations using python-pptx."""

    def can_process(self, file_path: str) -> bool:
        return Path(file_path).suffix.lower() in ('.pptx', '.ppt')

    def process(self, file_path: str, metadata: Optional[Dict[str, Any]] = None) -> Document:
        base_metadata = self._extract_metadata(file_path)
        if metadata:
            base_metadata.update(metadata)

        try:
            prs = Presentation(file_path)

            content_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                if slide_text:
                    content_parts.append(f"[Slide {slide_num}]\n" + '\n'.join(slide_text))

            content = '\n\n'.join(content_parts)

            base_metadata.update({
                'total_slides': len(prs.slides),
                'document_type': 'pptx',
            })

            return Document(
                id=self._generate_id(file_path),
                content=content,
                metadata=base_metadata,
                source_file=file_path,
            )

        except Exception as e:
            raise ValueError(f"Failed to process PPTX: {str(e)}")


class TextProcessor(DocumentProcessor):
    """Processor for plain text files."""

    def can_process(self, file_path: str) -> bool:
        return Path(file_path).suffix.lower() in ('.txt', '.md', '.csv', '.json', '.xml', '.html')

    def process(self, file_path: str, metadata: Optional[Dict[str, Any]] = None) -> Document:
        base_metadata = self._extract_metadata(file_path)
        if metadata:
            base_metadata.update(metadata)

        try:
            # Try different encodings
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            content = None

            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    break
                except UnicodeDecodeError:
                    continue

            if content is None:
                content = ''  # Empty file or encoding issues

            # Count lines, words, characters
            lines = content.split('\n')
            words = content.split()

            base_metadata.update({
                'total_lines': len(lines),
                'total_words': len(words),
                'total_characters': len(content),
                'document_type': 'text',
            })

            return Document(
                id=self._generate_id(file_path),
                content=content,
                metadata=base_metadata,
                source_file=file_path,
            )

        except Exception as e:
            raise ValueError(f"Failed to process text file: {str(e)}")


class DocumentProcessorFactory:
    """Factory for creating the appropriate document processor."""

    def __init__(self):
        self._processors: List[DocumentProcessor] = [
            PDFProcessor(),
            DocxProcessor(),
            PPTXProcessor(),
            TextProcessor(),
        ]

    def get_processor(self, file_path: str) -> Optional[DocumentProcessor]:
        """Get the appropriate processor for the given file."""
        for processor in self._processors:
            if processor.can_process(file_path):
                return processor
        return None

    def process(self, file_path: str, metadata: Optional[Dict[str, Any]] = None) -> Document:
        """Process a document using the appropriate processor."""
        processor = self.get_processor(file_path)
        if not processor:
            raise ValueError(f"Unsupported file type: {Path(file_path).suffix}")
        return processor.process(file_path, metadata)


# Singleton instance
_factory_instance = None


def get_document_processor() -> DocumentProcessorFactory:
    """Get the singleton document processor factory."""
    global _factory_instance
    if _factory_instance is None:
        _factory_instance = DocumentProcessorFactory()
    return _factory_instance
